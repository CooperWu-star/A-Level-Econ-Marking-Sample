"""Assemble a CIE Economics slide deck (.pptx) from a JSON spec.

Usage:
    python build_deck.py spec.json out.pptx
    python build_deck.py - out.pptx          # read spec from stdin

------------------------------------------------------------------------------
DESIGN SYSTEM (v2 — ported from the ppt-master design language)
------------------------------------------------------------------------------
The look is governed by three things, all baked into this engine so the spec
stays about *content*:

1. PALETTE — pick one with the top-level "palette" field (default "academy").
   Each palette is a role set {bg, surface, primary, accent, ink, muted}.
   The 60-30-10 rule is honoured: bg is the breathing field, primary is the
   structural colour, accent is used sparingly (rules, chips, one number).

2. LOCKED TYPE SCALE — every structural role resolves to ONE size, reused on
   every slide (same-role size drift is the #1 "AI deck" tell). See SCALE.

3. PAGE RHYTHM — slides are tagged anchor / dense / breathing and styled
   accordingly, so the deck does NOT read as "a coloured band + bullets on
   every page". Section dividers, big-number and quote slides punctuate the
   dense content runs.

------------------------------------------------------------------------------
SPEC SHAPE
------------------------------------------------------------------------------
{
  "course":   "CIE A-Level Economics 9708",     # or "IGCSE Economics 0455"
  "chapter":  "Topic 2.1 - Demand and supply curves",
  "subtitle": "AS Level | Microeconomics",       # optional
  "footer":   "Aixiom Academy",                  # optional, every slide
  "palette":  "academy",                          # optional, see PALETTES
  "kicker":   "9708 - 2.1 - Coursebook Unit 2, Ch 7",  # optional top-line tag
  "slides": [
    { "type": "title",   "title": "...", "subtitle": "..." },
    { "type": "section", "title": "Part 1 - Determinants", "number": "1",
      "subtitle": "optional one-line framing" },
    { "type": "definition", "title": "Effective demand", "term": "Effective demand",
      "definition": "...", "bullets": ["...", "..."], "notes": "..." },
    { "type": "content", "title": "Determinants of demand",
      "intro": "optional lead line", "bullets": ["...", "..."] },
    { "type": "diagram", "title": "Shift in demand",
      "diagram_id": "demand_supply_demand_right",
      "diagram_title": "optional override", "caption": "...",
      "bullets": ["...", "..."] },
    { "type": "case_study", "title": "Case: 2022-24 egg price spike (US)",
      "era": "Recent", "summary": "...", "bullets": ["..."],
      "link_to_theory": "Leftward shift in S -> higher P, lower Q." },
    { "type": "compare", "title": "Movement vs shift",
      "left":  {"head": "Movement along", "bullets": ["..."]},
      "right": {"head": "Shift of the curve", "bullets": ["..."]} },
    { "type": "stat", "title": "Scale of the shock",
      "value": "+60%", "label": "US egg prices, YoY 2022-23",
      "caption": "Low PED -> price does most of the adjusting." },
    { "type": "quote", "text": "Price is the most powerful signal in a market.",
      "attribution": "- exam framing" },
    { "type": "summary", "title": "Topic recap", "bullets": ["..."] }
  ]
}

Every slide may carry "notes" (speaker notes) and "rhythm" (anchor / dense /
breathing) to override the default for that slide type.
"""
from __future__ import annotations
import json, os, sys, tempfile, shutil, uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))
import diagrams as dg  # noqa: E402

# ---------------------------------------------------------------------------
# Geometry  (13.333 x 7.5in 16:9  ==  ppt-master's 1280x720 canvas / 96)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MX = Inches(0.62)                       # left / right margin  (~60px)
CONTENT_L = MX
CONTENT_W = Inches(13.333 - 2 * 0.62)   # 12.093"
TITLE_KICKER_Y = Inches(0.40)
TITLE_Y = Inches(0.62)
RULE_Y = Inches(1.42)
CONTENT_TOP = Inches(1.66)
CONTENT_BOTTOM = Inches(6.95)
FOOTER_Y = Inches(7.04)

# ---------------------------------------------------------------------------
# Palettes  (role set; hex adapted from ppt-master image-palettes)
#   bg      breathing field (~60%)
#   surface card / panel fill
#   primary structural colour (~30%)
#   accent  sparing pop (<10%): rules, chips, one number
#   ink     body text on bg/surface
#   muted   secondary text / hairlines
#   dark    True -> body text is light, surfaces are deep
# ---------------------------------------------------------------------------
def _rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


PALETTES = {
    # Default - keeps the Aixiom navy/orange identity, cleaned up
    "academy": dict(bg="FFFFFF", surface="F4F6F8", primary="0E3A66", accent="E06B2A",
                    ink="1F2A37", muted="6B7280", dark=False, head_font="Calibri"),
    "cool-corporate": dict(bg="F8F9FA", surface="FFFFFF", primary="1E3A5F", accent="C99A2E",
                           ink="1A2433", muted="6B7280", dark=False, head_font="Calibri"),
    "editorial-classic": dict(bg="FAF7F2", surface="FFFFFF", primary="0F2C4C", accent="C2410C",
                              ink="1A1A1A", muted="6B6257", dark=False, head_font="Georgia"),
    "warm-earth": dict(bg="FEF7ED", surface="FFFFFF", primary="9A3412", accent="B45309",
                       ink="2A1E16", muted="8A7563", dark=False, head_font="Georgia"),
    "nature-organic": dict(bg="F6FAF4", surface="FFFFFF", primary="166534", accent="B7791F",
                           ink="1B2A1B", muted="5F6B5F", dark=False, head_font="Calibri"),
    "mono-ink": dict(bg="FFFFFF", surface="F4F4F5", primary="1A1A1A", accent="C2410C",
                     ink="1A1A1A", muted="6B7280", dark=False, head_font="Arial"),
    "dark-cinematic": dict(bg="0B1020", surface="1E293B", primary="14B8A6", accent="D4AF37",
                           ink="E5E7EB", muted="94A3B8", dark=True, head_font="Calibri"),
}

# ---------------------------------------------------------------------------
# Locked type scale (pt) - one size per role, reused deck-wide.
# Tuned a little larger than ppt-master's "balanced" baseline because a
# teaching deck is projected and read from the back of a classroom.
# ---------------------------------------------------------------------------
SCALE = dict(
    cover=40, section=34, kpinum=54, term=28, pagetitle=26,
    subtitle=18, lead=17, body=16, small=13, caption=12, kicker=11, foot=9,
)
BODY_FONT = "Calibri"


class Theme:
    def __init__(self, name):
        p = PALETTES.get(name, PALETTES["academy"])
        self.bg = _rgb(p["bg"]); self.surface = _rgb(p["surface"])
        self.primary = _rgb(p["primary"]); self.accent = _rgb(p["accent"])
        self.ink = _rgb(p["ink"]); self.muted = _rgb(p["muted"])
        self.dark = p["dark"]; self.head_font = p["head_font"]
        # text colours that flip on dark themes
        self.on_primary = _rgb("FFFFFF")
        self.title_ink = self.primary if not self.dark else _rgb("F8FAFC")
        self.body_ink = self.ink
        # surface for cards always reads as a light panel even on dark decks,
        # so embedded matplotlib diagrams (white bg) never clash
        self.card = self.surface


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _noline_noshadow(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def _rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    _noline_noshadow(sp)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    return sp


def _text(slide, l, t, w, h, text, *, size, color, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None, spacing=1.08,
          letter_caps=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text.upper() if letter_caps else text
    f = run.font
    f.name = font or BODY_FONT
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return tb


def _bullets(slide, l, t, w, h, items, theme, *, size=None, gap=10, glyph="▪"):
    """Bullets with a hanging accent glyph.

    Each item is either a plain string (a cue) or a dict that TEACHES:
      {"point": "Income", "detail": "Normal good → D shifts right as income rises."}
    The point is bolded; the detail explains it on the same line. Use the dict
    form on teaching slides so the deck explains rather than just lists.
    """
    if not items:
        return None
    size = size or SCALE["body"]
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.14
        g = p.add_run(); g.text = f"{glyph}  "
        g.font.name = BODY_FONT; g.font.size = Pt(size)
        g.font.color.rgb = theme.accent; g.font.bold = True
        if isinstance(item, dict):
            point = item.get("point") or item.get("term") or ""
            detail = item.get("detail") or item.get("desc") or ""
            pr = p.add_run(); pr.text = point + (" — " if detail else "")
            pr.font.name = BODY_FONT; pr.font.size = Pt(size)
            pr.font.bold = True; pr.font.color.rgb = theme.body_ink
            if detail:
                dr = p.add_run(); dr.text = detail
                dr.font.name = BODY_FONT; dr.font.size = Pt(size)
                dr.font.color.rgb = theme.body_ink
        else:
            r = p.add_run(); r.text = item
            r.font.name = BODY_FONT; r.font.size = Pt(size)
            r.font.color.rgb = theme.body_ink
    return tb


def _hairline(slide, l, t, w, color, thick=Inches(0.012)):
    return _rect(slide, l, t, w, thick, color)


def _chip(slide, l, t, label, fill, text_color):
    w, h = Inches(1.25), Inches(0.34)
    chip = _rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        chip.adjustments[0] = 0.5
    except Exception:
        pass
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = BODY_FONT; r.font.size = Pt(SCALE["foot"] + 1)
    r.font.bold = True; r.font.color.rgb = text_color
    return chip


def _bg(slide, theme, color=None):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, color or theme.bg)


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Shared "dense page" chrome - LIGHT title treatment (kicker + title + rule).
# Replaces the old full-width coloured band that sat on every slide.
# ---------------------------------------------------------------------------
def _page_header(slide, theme, title, ctx, *, idx=None, total=None):
    _bg(slide, theme)
    # left accent tab - a small vertical mark, the only structural pop up top
    _rect(slide, CONTENT_L, TITLE_KICKER_Y + Inches(0.02), Inches(0.07), Inches(0.92), theme.accent)
    kicker = ctx.get("kicker") or f'{ctx["course"]}'
    _text(slide, CONTENT_L + Inches(0.22), TITLE_KICKER_Y, CONTENT_W, Inches(0.26),
          kicker, size=SCALE["kicker"], color=theme.muted, letter_caps=True)
    _text(slide, CONTENT_L + Inches(0.22), TITLE_Y, CONTENT_W, Inches(0.74),
          title, size=SCALE["pagetitle"], bold=True, color=theme.title_ink,
          font=theme.head_font)
    _hairline(slide, CONTENT_L, RULE_Y, CONTENT_W, theme.muted, thick=Inches(0.008))
    _footer(slide, theme, ctx, idx, total)


def _footer(slide, theme, ctx, idx=None, total=None):
    if ctx.get("footer"):
        _text(slide, CONTENT_L, FOOTER_Y, Inches(8), Inches(0.3),
              ctx["footer"], size=SCALE["foot"], color=theme.muted, italic=True)
    if idx is not None and total is not None:
        _text(slide, Inches(11.0), FOOTER_Y, Inches(1.71), Inches(0.3),
              f"{idx} / {total}", size=SCALE["foot"], color=theme.muted,
              align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def _slide_title(prs, s, ctx, theme):
    slide = _blank(prs)
    _bg(slide, theme, theme.primary)
    # accent rule
    _rect(slide, Inches(0.9), Inches(3.5), Inches(0.9), Inches(0.10), theme.accent)
    _text(slide, Inches(0.9), Inches(2.25), Inches(11.5), Inches(0.5),
          ctx["course"], size=SCALE["subtitle"], color=_tint(theme.on_primary),
          letter_caps=True)
    _text(slide, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.4),
          s.get("title", ctx["chapter"]), size=SCALE["cover"], bold=True,
          color=theme.on_primary, font=theme.head_font, spacing=1.02)
    sub = s.get("subtitle") or ctx.get("subtitle", "")
    if sub:
        _text(slide, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.6),
              sub, size=SCALE["subtitle"], color=_tint(theme.on_primary))
    if ctx.get("footer"):
        _text(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
              ctx["footer"], size=SCALE["caption"], color=_tint(theme.on_primary),
              italic=True)
    _notes(slide, s.get("notes"))


def _slide_section(prs, s, ctx, theme, idx, total):
    """Breathing divider - big number + title on a near-empty field."""
    slide = _blank(prs)
    _bg(slide, theme, theme.surface if not theme.dark else theme.primary)
    num = str(s.get("number", "")).strip()
    if num:
        _text(slide, CONTENT_L, Inches(1.9), Inches(4), Inches(2.2),
              num, size=120, bold=True, color=_soft(theme.accent),
              font=theme.head_font)
    _rect(slide, CONTENT_L + Inches(0.04), Inches(4.05), Inches(1.1), Inches(0.10), theme.accent)
    _text(slide, CONTENT_L, Inches(4.25), Inches(11.8), Inches(1.2),
          s.get("title", "Section"), size=SCALE["section"], bold=True,
          color=theme.title_ink, font=theme.head_font, spacing=1.0)
    if s.get("subtitle"):
        _text(slide, CONTENT_L, Inches(5.35), Inches(11.0), Inches(0.7),
              s["subtitle"], size=SCALE["lead"], color=theme.muted)
    _footer(slide, theme, ctx, idx, total)
    _notes(slide, s.get("notes"))


def _slide_definition(prs, s, ctx, theme, idx, total):
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Definition"), ctx, idx=idx, total=total)
    term = s.get("term") or s.get("title", "")
    _text(slide, CONTENT_L, CONTENT_TOP, CONTENT_W, Inches(0.7),
          term, size=SCALE["term"], bold=True, color=theme.primary,
          font=theme.head_font)
    # CIE definition in a quoted accent-edged card
    defn = s.get("definition", "")
    card_t = Inches(2.5)
    card = _rect(slide, CONTENT_L, card_t, CONTENT_W, Inches(1.5), theme.surface)
    _rect(slide, CONTENT_L, card_t, Inches(0.07), Inches(1.5), theme.accent)
    _text(slide, CONTENT_L + Inches(0.3), card_t + Inches(0.14), CONTENT_W - Inches(0.5),
          Inches(1.25), defn, size=SCALE["lead"], color=theme.body_ink,
          anchor=MSO_ANCHOR.MIDDLE, spacing=1.18)
    if s.get("bullets"):
        _text(slide, CONTENT_L, Inches(4.35), CONTENT_W, Inches(0.32),
              "Key points", size=SCALE["small"], bold=True, color=theme.muted,
              letter_caps=True)
        _bullets(slide, CONTENT_L, Inches(4.75), CONTENT_W, Inches(2.1),
                 s["bullets"], theme)
    _notes(slide, s.get("notes"))


def _slide_content(prs, s, ctx, theme, idx, total):
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Content"), ctx, idx=idx, total=total)
    top = CONTENT_TOP
    if s.get("intro"):
        _text(slide, CONTENT_L, top, CONTENT_W, Inches(0.8),
              s["intro"], size=SCALE["lead"], color=theme.muted, spacing=1.18)
        top = Inches(2.55)
    _bullets(slide, CONTENT_L, top, CONTENT_W, Inches(4.2),
             s.get("bullets", []), theme)
    _notes(slide, s.get("notes"))


def _slide_diagram(prs, s, ctx, theme, idx, total, tmpdir):
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Diagram"), ctx, idx=idx, total=total)
    png = os.path.join(tmpdir, f"{uuid.uuid4().hex}.png")
    dg.render(s["diagram_id"], png, title=s.get("diagram_title"))
    # diagram sits in a white card on the left (readable on any palette)
    card_l, card_t = CONTENT_L, CONTENT_TOP + Inches(0.05)
    card_w, card_h = Inches(6.5), Inches(4.9)
    _rect(slide, card_l, card_t, card_w, card_h, _rgb("FFFFFF"))
    slide.shapes.add_picture(png, card_l + Inches(0.18), card_t + Inches(0.2),
                             height=card_h - Inches(0.45))
    # right column: caption (lead) + walkthrough bullets
    rx = card_l + card_w + Inches(0.45)
    rw = CONTENT_L + CONTENT_W - rx
    ry = CONTENT_TOP + Inches(0.05)
    if s.get("caption"):
        _text(slide, rx, ry, rw, Inches(1.1),
              s["caption"], size=SCALE["lead"], bold=True, color=theme.primary,
              font=theme.head_font, spacing=1.16)
        ry = ry + Inches(1.2)
    if s.get("bullets"):
        _bullets(slide, rx, ry, rw, CONTENT_BOTTOM - ry, s["bullets"], theme,
                 size=SCALE["body"])
    _notes(slide, s.get("notes"))


def _slide_diagram_pair(prs, s, ctx, theme, idx, total, tmpdir):
    """Two diagrams side by side — for visual contrasts (movement vs shift,
    monopoly vs perfect competition, demand-pull vs cost-push, …)."""
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Comparison"), ctx, idx=idx, total=total)
    gap = Inches(0.4)
    col_w = (CONTENT_W - gap) / 2
    top = CONTENT_TOP + Inches(0.05)
    card_h = Inches(4.25)
    for k, x in (("left", CONTENT_L), ("right", CONTENT_L + col_w + gap)):
        blk = s.get(k) or {}
        _rect(slide, x, top, col_w, card_h + Inches(0.55), theme.surface)
        did = blk.get("diagram_id")
        if did:
            png = os.path.join(tmpdir, f"{uuid.uuid4().hex}.png")
            dg.render(did, png, title=blk.get("diagram_title"))
            pic = slide.shapes.add_picture(png, x, top + Inches(0.18),
                                           height=card_h - Inches(0.35))
            pic.left = int(x + (col_w - pic.width) / 2)
        if blk.get("caption"):
            _text(slide, x + Inches(0.2), top + card_h + Inches(0.05),
                  col_w - Inches(0.4), Inches(0.5), blk["caption"],
                  size=SCALE["small"], bold=True, color=theme.primary,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  font=theme.head_font)
    _notes(slide, s.get("notes"))


def _slide_case_study(prs, s, ctx, theme, idx, total):
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Case study"), ctx, idx=idx, total=total)
    era = (s.get("era") or "").strip().lower()
    if era.startswith("class"):
        _chip(slide, CONTENT_L, CONTENT_TOP, "CLASSIC", theme.primary, theme.on_primary)
    elif era.startswith(("recent", "modern")):
        _chip(slide, CONTENT_L, CONTENT_TOP, "RECENT", theme.accent, _rgb("FFFFFF"))
    top = CONTENT_TOP + Inches(0.5)
    if s.get("summary"):
        _text(slide, CONTENT_L, top, CONTENT_W, Inches(1.05),
              s["summary"], size=SCALE["lead"], color=theme.body_ink, spacing=1.18)
        top = top + Inches(1.15)
    if s.get("bullets"):
        _text(slide, CONTENT_L, top, CONTENT_W, Inches(0.32),
              "What it shows", size=SCALE["small"], bold=True, color=theme.muted,
              letter_caps=True)
        _bullets(slide, CONTENT_L, top + Inches(0.4), CONTENT_W, Inches(2.0),
                 s["bullets"], theme)
    if s.get("link_to_theory"):
        bt = Inches(6.05)
        _rect(slide, CONTENT_L, bt, CONTENT_W, Inches(0.78), theme.surface)
        _rect(slide, CONTENT_L, bt, Inches(0.07), Inches(0.78), theme.accent)
        _text(slide, CONTENT_L + Inches(0.3), bt, CONTENT_W - Inches(0.5), Inches(0.78),
              f"Theory link  -  {s['link_to_theory']}", size=SCALE["small"],
              italic=True, color=theme.primary, anchor=MSO_ANCHOR.MIDDLE)
    _notes(slide, s.get("notes"))


def _slide_compare(prs, s, ctx, theme, idx, total):
    """Two-column comparison - identical framing IS the message."""
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Comparison"), ctx, idx=idx, total=total)
    gap = Inches(0.4)
    col_w = (CONTENT_W - gap) / 2
    top = CONTENT_TOP + Inches(0.05)
    col_h = Inches(4.9)
    for k, x in (("left", CONTENT_L), ("right", CONTENT_L + col_w + gap)):
        block = s.get(k) or {}
        _rect(slide, x, top, col_w, col_h, theme.surface)
        _rect(slide, x, top, col_w, Inches(0.62), theme.primary)
        _text(slide, x + Inches(0.22), top, col_w - Inches(0.4), Inches(0.62),
              block.get("head", ""), size=SCALE["subtitle"], bold=True,
              color=theme.on_primary, anchor=MSO_ANCHOR.MIDDLE, font=theme.head_font)
        _bullets(slide, x + Inches(0.22), top + Inches(0.82), col_w - Inches(0.44),
                 col_h - Inches(1.0), block.get("bullets", []), theme,
                 size=SCALE["body"])
    _notes(slide, s.get("notes"))


def _slide_stat(prs, s, ctx, theme, idx, total):
    """Breathing big-number slide."""
    slide = _blank(prs)
    _bg(slide, theme, theme.surface if not theme.dark else theme.primary)
    _rect(slide, CONTENT_L, Inches(1.0), Inches(0.07), Inches(0.7), theme.accent)
    _text(slide, CONTENT_L + Inches(0.22), Inches(1.0), CONTENT_W, Inches(0.7),
          s.get("title", ""), size=SCALE["pagetitle"], bold=True,
          color=theme.title_ink, font=theme.head_font)
    _text(slide, CONTENT_L, Inches(2.5), CONTENT_W, Inches(1.7),
          s.get("value", ""), size=SCALE["kpinum"] + 30, bold=True,
          color=theme.accent, font=theme.head_font)
    if s.get("label"):
        _text(slide, CONTENT_L, Inches(4.45), CONTENT_W, Inches(0.6),
              s["label"], size=SCALE["subtitle"], bold=True, color=theme.title_ink)
    if s.get("caption"):
        _text(slide, CONTENT_L, Inches(5.15), Inches(10.5), Inches(0.9),
              s["caption"], size=SCALE["lead"], color=theme.muted, spacing=1.18)
    _footer(slide, theme, ctx, idx, total)
    _notes(slide, s.get("notes"))


def _slide_quote(prs, s, ctx, theme, idx, total):
    slide = _blank(prs)
    _bg(slide, theme, theme.surface if not theme.dark else theme.primary)
    _text(slide, CONTENT_L, Inches(1.4), Inches(2), Inches(1.4),
          "“", size=140, bold=True, color=_soft(theme.accent),
          font=theme.head_font)
    _text(slide, CONTENT_L + Inches(0.1), Inches(2.7), Inches(11.4), Inches(2.6),
          s.get("text", ""), size=28, color=theme.title_ink,
          font=theme.head_font, spacing=1.2)
    if s.get("attribution"):
        _text(slide, CONTENT_L + Inches(0.1), Inches(5.5), Inches(11.0), Inches(0.5),
              s["attribution"], size=SCALE["subtitle"], italic=True, color=theme.muted)
    _footer(slide, theme, ctx, idx, total)
    _notes(slide, s.get("notes"))


def _slide_summary(prs, s, ctx, theme, idx, total):
    slide = _blank(prs)
    _page_header(slide, theme, s.get("title", "Topic recap"), ctx, idx=idx, total=total)
    _text(slide, CONTENT_L, CONTENT_TOP, CONTENT_W, Inches(0.45),
          "You should now be able to…", size=SCALE["subtitle"], bold=True,
          color=theme.primary, font=theme.head_font)
    _bullets(slide, CONTENT_L, CONTENT_TOP + Inches(0.6), CONTENT_W, Inches(4.0),
             s.get("bullets", []), theme, glyph="✓")
    _notes(slide, s.get("notes"))


# ---------- small colour helpers ----------
def _tint(color):
    """Lighten toward a soft on-primary subtitle colour."""
    return RGBColor(min(color[0], 220), min(color[1], 222), min(color[2], 228)) \
        if color == RGBColor(255, 255, 255) else RGBColor(0xCC, 0xD6, 0xE0)


def _soft(color):
    """A pale wash of an accent for oversized decorative glyphs/numbers."""
    r, g, b = color
    mix = lambda c: int(c + (255 - c) * 0.72)
    return RGBColor(mix(r), mix(g), mix(b))


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------
def build(spec: dict, out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    theme = Theme(spec.get("palette", "academy"))
    # tie the diagram accent to the deck accent (shifts / new curves)
    try:
        dg.set_theme(accent="#%02X%02X%02X" % (theme.accent[0], theme.accent[1], theme.accent[2]))
    except Exception:
        pass

    ctx = {
        "course": spec.get("course", "CIE Economics"),
        "chapter": spec.get("chapter", ""),
        "subtitle": spec.get("subtitle", ""),
        "footer": spec.get("footer", ""),
        "kicker": spec.get("kicker", ""),
    }

    slides = spec.get("slides", [])
    content_slides = [s for s in slides if s.get("type") != "title"]
    total = len(content_slides)

    tmpdir = tempfile.mkdtemp(prefix="ciepptx_")
    try:
        idx = 0
        for s in slides:
            stype = s.get("type", "content")
            if stype == "title":
                _slide_title(prs, s, ctx, theme)
                continue
            idx += 1
            if stype == "section":
                _slide_section(prs, s, ctx, theme, idx, total)
            elif stype == "definition":
                _slide_definition(prs, s, ctx, theme, idx, total)
            elif stype == "diagram":
                _slide_diagram(prs, s, ctx, theme, idx, total, tmpdir)
            elif stype == "diagram_pair":
                _slide_diagram_pair(prs, s, ctx, theme, idx, total, tmpdir)
            elif stype == "case_study":
                _slide_case_study(prs, s, ctx, theme, idx, total)
            elif stype == "compare":
                _slide_compare(prs, s, ctx, theme, idx, total)
            elif stype == "stat":
                _slide_stat(prs, s, ctx, theme, idx, total)
            elif stype == "quote":
                _slide_quote(prs, s, ctx, theme, idx, total)
            elif stype == "summary":
                _slide_summary(prs, s, ctx, theme, idx, total)
            else:
                _slide_content(prs, s, ctx, theme, idx, total)
        prs.save(out_path)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return out_path


def main(argv):
    if len(argv) < 3:
        print("Usage: build_deck.py <spec.json|-> <out.pptx>", file=sys.stderr)
        return 2
    spec_arg, out_path = argv[1], argv[2]
    if spec_arg == "-":
        spec = json.load(sys.stdin)
    else:
        with open(spec_arg, "r", encoding="utf-8") as f:
            spec = json.load(f)
    build(spec, out_path)
    print(out_path)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
